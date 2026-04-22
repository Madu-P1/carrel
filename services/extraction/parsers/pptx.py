from __future__ import annotations

from pathlib import Path

from fastapi import HTTPException

from ..types import ExtractedElement
from ..utils import normalize_space, file_sha
from .common import ParserContext, build_asset, make_span

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None


def parse_pptx(path: Path, *, suffix: str, mime_type: str, context: ParserContext):
    if PptxPresentation is None:
        raise HTTPException(status_code=400, detail="PPTX support requires python-pptx")
    prs = PptxPresentation(str(path))
    file_id = file_sha(path)[:16]
    elements: list[ExtractedElement] = []
    warnings: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_title = None
        if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
            slide_title = normalize_space(slide.shapes.title.text)
            if slide_title:
                span = make_span(path, file_id, slide=slide_index, section=slide_title, element_id=f"slide-{slide_index}-title")
                elements.append(
                    ExtractedElement(
                        id=span.element_id or f"slide-{slide_index}-title",
                        kind="slide",
                        text=slide_title,
                        normalized_text=slide_title,
                        span=span,
                        metadata={"slide_number": slide_index},
                    )
                )
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    text = normalize_space(paragraph.text)
                    if not text:
                        continue
                    span = make_span(
                        path,
                        file_id,
                        slide=slide_index,
                        section=slide_title or f"Slide {slide_index}",
                        paragraph_id=f"s{slide_index}-p{shape_index}-{paragraph_index}",
                        element_id=f"pptx-{slide_index}-{shape_index}-{paragraph_index}",
                    )
                    elements.append(
                        ExtractedElement(
                            id=span.element_id or f"pptx-{slide_index}-{shape_index}-{paragraph_index}",
                            kind="bullet_list" if paragraph.level > 0 else "paragraph",
                            text=text,
                            normalized_text=text,
                            span=span,
                            metadata={"slide_number": slide_index, "paragraph_level": paragraph.level},
                        )
                    )
            if getattr(shape, "has_table", False):
                rows: list[str] = []
                for row in shape.table.rows:
                    values = [normalize_space(cell.text) for cell in row.cells]
                    if any(values):
                        rows.append(" | ".join(values))
                if rows:
                    table_text = "\n".join(rows)
                    span = make_span(
                        path,
                        file_id,
                        slide=slide_index,
                        section=slide_title or f"Slide {slide_index}",
                        element_id=f"pptx-table-{slide_index}-{shape_index}",
                    )
                    elements.append(
                        ExtractedElement(
                            id=span.element_id or f"pptx-table-{slide_index}-{shape_index}",
                            kind="table",
                            text=table_text,
                            normalized_text=table_text,
                            span=span,
                            metadata={"slide_number": slide_index},
                        )
                    )
        try:
            notes_slide = slide.notes_slide
            note_lines = []
            for shape in notes_slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    note_text = normalize_space(shape.text)
                    if note_text and "click to add notes" not in note_text.lower():
                        note_lines.append(note_text)
            if note_lines:
                joined = "\n".join(note_lines)
                span = make_span(path, file_id, slide=slide_index, section=slide_title or f"Slide {slide_index}", element_id=f"pptx-notes-{slide_index}")
                elements.append(
                    ExtractedElement(
                        id=span.element_id or f"pptx-notes-{slide_index}",
                        kind="speaker_note",
                        text=joined,
                        normalized_text=joined,
                        span=span,
                        metadata={"slide_number": slide_index},
                    )
                )
        except Exception:
            warnings.append(f"Speaker notes could not be loaded for slide {slide_index}.")
    return build_asset(
        path,
        detected_type=suffix,
        mime_type=mime_type,
        parser_name="python-pptx-structured",
        elements=elements,
        context=context,
        warnings=warnings,
        extraction_modes=["slides", "speaker_notes"],
        metadata={"page_count": len(prs.slides)},
        confidence=0.83,
    )
