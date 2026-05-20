"""Integration coverage for the five Docling formats added in T10.

T10 (Phase 4.1) extends `INGEST_DOCLING_FORMATS` beyond `pdf` so DOCX,
HTML, Markdown, PPTX, and LaTeX also route through the Docling
typed-node ingest path. This suite proves each newly-supported format
parses into typed nodes end to end:

    parse_document -> DoclingDocument -> typed_walker.walk -> TypedNode[]

`DoclingFormatCoverageTests` is skipped when Docling is not installed
(it is an optional ~1-2 GB dependency, gated by
`docling_parser.is_available()`), matching the optional-dependency
design of `services.ingestion.docling_parser`. PDF coverage already
ships via `evals/fixtures/cell_division.pdf` (T58); this suite covers
only the five formats T10 adds.

`DoclingDefaultAllowlistTests` needs no Docling install: it pins the
`INGEST_DOCLING_FORMATS` default that T10 widened.
"""

from __future__ import annotations

import os
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from services.ingestion import docling_parser, typed_walker
from services.ingestion.orchestrator import _docling_enabled_for


def _write_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_heading("Mitosis", level=1)
    document.add_paragraph(
        "Mitosis is the process by which a eukaryotic cell divides into "
        "two genetically identical daughter cells."
    )
    document.save(str(path))


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Cell Division"
    slide.placeholders[1].text = "Meiosis produces four haploid gametes from a single diploid cell."
    presentation.save(str(path))


@unittest.skipUnless(
    docling_parser.is_available(),
    "Docling is not installed (optional ~1-2 GB dependency).",
)
class DoclingFormatCoverageTests(unittest.TestCase):
    """One case per format T10 adds to the Docling allowlist."""

    @classmethod
    def setUpClass(cls) -> None:
        cls._tmpdir = tempfile.TemporaryDirectory(prefix="t10-docling-")
        root = Path(cls._tmpdir.name)

        _write_docx(root / "mitosis.docx")
        _write_pptx(root / "cell-division.pptx")
        (root / "photosynthesis.html").write_text(
            "<html><body><h1>Photosynthesis</h1>"
            "<p>Photosynthesis converts light energy into chemical energy "
            "stored in glucose.</p></body></html>",
            encoding="utf-8",
        )
        (root / "osmosis.md").write_text(
            "# Osmosis\n\n"
            "Osmosis is the diffusion of water across a semipermeable "
            "membrane toward a higher solute concentration.\n",
            encoding="utf-8",
        )
        (root / "diffusion.tex").write_text(
            "\\documentclass{article}\n"
            "\\begin{document}\n"
            "\\section{Diffusion}\n"
            "Diffusion is the net movement of particles from a region of "
            "high concentration to one of low concentration.\n"
            "\\end{document}\n",
            encoding="utf-8",
        )
        cls._root = root

    @classmethod
    def tearDownClass(cls) -> None:
        cls._tmpdir.cleanup()

    def _assert_parses_to_nodes(self, filename: str, keyword: str) -> None:
        """parse_document + typed_walker.walk yields >=1 typed node that
        carries the fixture's distinctive keyword."""
        document = docling_parser.parse_document(self._root / filename)
        nodes = typed_walker.walk(document)

        self.assertGreaterEqual(len(nodes), 1, msg=f"{filename}: Docling produced zero typed nodes")
        combined = " ".join(node.verbatim_text for node in nodes).lower()
        self.assertIn(
            keyword,
            combined,
            msg=f"{filename}: expected '{keyword}' in extracted text, got {combined!r}",
        )

    def test_docx_parses_to_typed_nodes(self) -> None:
        self._assert_parses_to_nodes("mitosis.docx", "mitosis")

    def test_html_parses_to_typed_nodes(self) -> None:
        self._assert_parses_to_nodes("photosynthesis.html", "photosynthesis")

    def test_md_parses_to_typed_nodes(self) -> None:
        self._assert_parses_to_nodes("osmosis.md", "osmosis")

    def test_pptx_parses_to_typed_nodes(self) -> None:
        self._assert_parses_to_nodes("cell-division.pptx", "meiosis")

    def test_latex_parses_to_typed_nodes(self) -> None:
        self._assert_parses_to_nodes("diffusion.tex", "diffusion")


class DoclingDefaultAllowlistTests(unittest.TestCase):
    """`INGEST_DOCLING_FORMATS` default widened by T10. No Docling needed."""

    def test_default_allowlist_covers_six_formats_and_excludes_epub_txt(self) -> None:
        with mock.patch.dict("os.environ", {"INGEST_USE_DOCLING": "true"}, clear=False):
            os.environ.pop("INGEST_DOCLING_FORMATS", None)
            # `tex` is the file extension for InputFormat.LATEX; the
            # allowlist is extension-keyed, not InputFormat-name-keyed.
            for extension in ("pdf", "docx", "html", "md", "pptx", "tex"):
                self.assertTrue(
                    _docling_enabled_for(extension),
                    msg=f"{extension} should be in the default Docling allowlist",
                )
            # EPUB and TXT are not Docling InputFormat members; they stay
            # on the legacy extraction path (operator decision 2026-05-20).
            for extension in ("epub", "txt"):
                self.assertFalse(
                    _docling_enabled_for(extension),
                    msg=f"{extension} must not be in the default Docling allowlist",
                )


if __name__ == "__main__":
    unittest.main()
